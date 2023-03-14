'''
Convert a "states" directory to a PowerPoint presentation
'''
from ase.io import read
from ase.io import write
from ase.neb import NEB
from pptx import Presentation
from pptx.util import Inches

def find_states(statespath):
    '''
    Find the states in the AKMC simulation

    Statespath is the path to the states directory. In this
    directory there is a file "state_table". This file
    lists all the states and their corresponding energy.
    We just read the state label and return them in a list.
    '''
    statetable = statespath+"/state_table"
    fp = open(statetable)
    intable = fp.readlines()
    fp.close()
    table=[]
    for state in intable:
        state = state.split()
        table.append(state[0])
    return table

def find_processes(statespath,statename):
    '''
    Find the processes for a given state
    '''
    processtable = statespath+"/"+statename+"/processtable"
    fp = open(processtable)
    intable = fp.readlines()
    intable = intable[1:] # The first line is a comment
    table = []
    for process in intable:
        process = process.split()
        if process[3] != "-1":
            table.append(process[0])
    return table

def convert_con2img(statespath,statename,processname,movies):
    '''
    Convert the configuration files to image files

    For every process there 3 files:
    - reactant
    - saddle point
    - product
    we convert all these files to image files.
    '''
    fn_reactant = statespath+"/"+statename+"/procdata/reactant_"+processname
    fn_saddle   = statespath+"/"+statename+"/procdata/saddle_"+processname
    fn_product  = statespath+"/"+statename+"/procdata/product_"+processname
    fn_reaction = statespath+"/"+statename+"/procdata/reaction_"+processname
    structure_r = read(fn_reactant+".con")
    structure_s = read(fn_saddle+".con")
    structure_p = read(fn_product+".con")
    if movies:
        images = [structure_r]
        images += [structure_r.copy() for i in range(8)]
        images += [structure_p]
        neb = NEB(images)
        neb.interpolate()
        write(fn_reaction+".gif",images)
    else:
        write(fn_reactant+".png",structure_r)
        write(fn_saddle+".png",structure_s)
        write(fn_product+".png",structure_p)
    

def pptx_process(prs,statespath,statename,processname,movies):
    '''
    Add the process images to the PowerPoint file

    We put all images into a table. Here we add one row to the table.
    '''
    fn_reactant = statespath+"/"+statename+"/procdata/reactant_"+processname+".png"
    fn_saddle   = statespath+"/"+statename+"/procdata/saddle_"+processname+".png"
    fn_product  = statespath+"/"+statename+"/procdata/product_"+processname+".png"
    fn_reaction = statespath+"/"+statename+"/procdata/reaction_"+processname+".gif"
    layout = prs.slide_layouts[5]
    slide  = prs.slides.add_slide(layout)
    width  = prs.slide_width
    title  = slide.shapes.title
    title.text = "State: "+str(statename)+" Process: "+str(processname)
    if movies:
        top   = Inches(2.0)
        left  = Inches(1.0)
        width = Inches(5.5)
        height= Inches(5.5)
        img   = slide.shapes.add_picture(fn_reaction,left,top,width,height)
    else:
        width = Inches(3.0)
        height= Inches(3.0)
        top   = Inches(2.0)
        left1 = Inches(1.0)
        left2 = Inches(5.0)
        left3 = Inches(9.0)
        img   = slide.shapes.add_picture(fn_reactant,left1,top,width,height)
        img   = slide.shapes.add_picture(fn_saddle,left2,top,width,height)
        img   = slide.shapes.add_picture(fn_product,left3,top,width,height)

def doit(statespath,movies,template):
    '''
    Go through the whole thing
    '''
    if template:
        prs = Presentation(template)
    else:
        prs = Presentation()
        # Setup for widescreen slides
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    stateslist = find_states(statespath)
    for state in stateslist:
        processlist = find_processes(statespath,state)
        for process in processlist:
            convert_con2img(statespath,state,process,movies)
            pptx_process(prs,statespath,state,process,movies)
    prs.save("states.pptx")

def arguments():
    '''
    Read command line arguments.

    In this case the only available command line argument is
    the states pathname.
    '''
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("pathname", type=str, help="the pathname for the states directory, i.e. where the file \"state_table\" lives")
    parser.add_argument("--movies", dest="movies", default=False, action='store_const', const=True, help="if provided generate movies of the reaction processes")
    parser.add_argument("--template", dest="template", help="the PowerPoint template to start from")
    return parser.parse_args()

if __name__ == "__main__":
    args = arguments()
    doit(args.pathname,args.movies,args.template)
